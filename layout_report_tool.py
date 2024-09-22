from pptx import Presentation
import pandas as pd


def list_placeholders(design_number: int) -> pd.DataFrame:
    """
    Lists the placeholders in a PowerPoint design and returns a DataFrame.

    Parameters:
    design_number (int): The design number of the PowerPoint file.

    Returns:
    pd.DataFrame: DataFrame containing placeholder details.
    """
    placeholders_data = []

    # Load the presentation
    prs = Presentation(f"Powerpointer-main/Designs/Design-{design_number}.pptx")

    # Iterate over slide layouts and placeholders
    for i, layout in enumerate(prs.slide_layouts):
        for shape in layout.placeholders:
            placeholder_format = shape.placeholder_format
            placeholders_data.append([
                i,
                placeholder_format.idx,
                shape.shape_id,
                shape.name,
                placeholder_format.type,
                shape.width,
                shape.height,
                shape.left,
                shape.top
            ])

    # Create DataFrame
    df = pd.DataFrame(placeholders_data, columns=[
        'Layout Index', 'Placeholder Index', 'Shape ID', 'Name', 'Type', 'Width', 'Height', 'Left', 'Top'
    ])

    return df


def color_rows_by_layout(df: pd.DataFrame) -> pd.io.formats.style.Styler:
    """
    Applies color to rows of DataFrame based on the layout index.

    Parameters:
    df (pd.DataFrame): DataFrame containing placeholder details.

    Returns:
    pd.io.formats.style.Styler: Styler object with applied row colors.
    """

    def apply_colors(row):
        colors = [
            'background-color: #ffcccc',  # Light red
            'background-color: #ccffcc',  # Light green
            'background-color: #ccccff',  # Light blue
            'background-color: #ffffcc',  # Light yellow
            'background-color: #ffccff',  # Light pink
            'background-color: #ccffff',  # Light cyan
            'background-color: #ffd700',  # Gold
            'background-color: #e6e6fa'  # Lavender
        ]
        return [colors[row['Layout Index'] % len(colors)]] * len(row)

    return df.style.apply(apply_colors, axis=1)


def get_placeholder_indices_by_layout(df: pd.DataFrame) -> tuple:
    """
    Retrieves placeholder indices grouped by layout index and indices containing content placeholders.

    Parameters:
    df (pd.DataFrame): DataFrame containing placeholder details.

    Returns:
    tuple: A tuple containing:
        - List of lists of placeholder indices by layout.
        - List of unique layout indices.
        - List of lists of content placeholder indices by layout.
    """
    # Group by layout index and get placeholder indices
    placeholder_indices_by_layout = df.groupby('Layout Index')['Placeholder Index'].apply(list).tolist()
    layout_indices = df['Layout Index'].unique().tolist()

    # Filter and group DataFrame by content placeholders
    filtered_df = df[df['Name'].str.startswith('Content Placeholder')]
    grouped = filtered_df.groupby('Layout Index')['Placeholder Index'].apply(list)
    index_containing_placeholders = grouped.tolist()

    return placeholder_indices_by_layout, layout_indices, index_containing_placeholders


def supporting_parameters(design_number: int) -> tuple:
    """
    Generates supporting parameters for the given design number.

    Parameters:
    design_number (int): The design number of the PowerPoint file.

    Returns:
    tuple: A tuple containing:
        - List of lists of placeholder indices by layout.
        - List of unique layout indices.
        - List of lists of content placeholder indices by layout.
    """
    placeholders_df = list_placeholders(design_number)

    # Display the styled DataFrame (optional but very usefull if you want to check yur custom design details)
    styled_df = color_rows_by_layout(placeholders_df)

    # Get placeholder indices by layout
    placeholder_indices_by_layout_, layout_indices_, index_containing_placeholders = get_placeholder_indices_by_layout(
        placeholders_df)

    return placeholder_indices_by_layout_, layout_indices_, index_containing_placeholders
