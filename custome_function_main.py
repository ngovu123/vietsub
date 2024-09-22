import os
import random
import re

import google.generativeai as genai
from dotenv import load_dotenv
from pptx import Presentation

from Cache.default_prompt import prompt
from content_extractor import extract_contents_from_text
from layout_report_tool import supporting_parameters

# Load environment variables from .env file
load_dotenv()

# Get the API key from the environment
api_key = os.getenv("API_KEY")

# Configure the genai library with the API key
genai.configure(api_key=api_key)


def get_bot_response(topic: str, theme: str) -> tuple:
    """
    Generates a PowerPoint presentation based on the provided topic and theme.

    Parameters:
    topic (str): The topic for the presentation.
    theme (str): The theme for the presentation.

    Returns:
    tuple: A tuple containing the path to the generated PowerPoint file and the filename.
    """
    user_text = topic
    text = ""
    input_string = re.sub(r'[^\w\s.\-\(\)]', '', user_text).replace("\n", "")
    number = int(theme[-1])
    pptlink = None

    if number not in range(1, 10):
        number = 1
        print("Unavailable design, using default design...")

    print(f"Available design, using {theme} design...")

    # Generate a filename using OpenAI API
    filename_prompt = f"""Generate a short, descriptive filename based on the following input: \"{user_text}\".
Answer just with the short filename; no other explanation. 
Do not give extensions to files like my_file.txt. I just need a file name."""

    model = genai.GenerativeModel('gemini-pro')

    generation_config = genai.GenerationConfig(
        stop_sequences=None,
        temperature=0.9,
        top_p=1.0,
        top_k=32,
        candidate_count=1,
        max_output_tokens=400,
    )

    filename_response = model.generate_content(
        contents=filename_prompt,
        generation_config=generation_config,
        stream=False,
    )

    filename = filename_response.text.strip().replace(" ", "_")

    cache_dir = 'Powerpointer-main/Cache'
    if not os.path.exists(cache_dir):
        os.makedirs(cache_dir)

    if number <= 7:
        question = prompt(user_text)

        text = model.generate_content(
            contents=question,
            generation_config=generation_config,
            stream=False,
        )

        with open(f'{cache_dir}/{filename}.txt', 'w', encoding='utf-8') as f:
            f.write(text.text)

        placeholder_indices_by_layout_, layout_indices_, _ = supporting_parameters(number)
        pptlink = create_ppt_default(f'{cache_dir}/{filename}.txt', number, filename)

    else:
        _, _, index_containing_placeholders = supporting_parameters(number)
        pptlink = create_ppt_custom(f'Powerpointer-main/Cache/custom_prompt.txt', number, filename,
                                    index_containing_placeholders)

    return pptlink, f'{cache_dir}/{filename}'


def create_ppt_custom(text_file: str, design_number: int, ppt_name: str, index_containing_placeholders: list) -> str:
    """
    Creates a custom PowerPoint presentation based on the provided text file and design number.

    Parameters:
    text_file (str): Path to the text file containing slide contents.
    design_number (int): The design number of the PowerPoint template.
    ppt_name (str): The name for the generated PowerPoint file.
    index_containing_placeholders (list): List of indices containing placeholders.

    Returns:
    str: Path to the generated PowerPoint file.
    """
    prs = Presentation(f"Powerpointer-main/Designs/Design-{design_number}.pptx")
    slide_count = 0
    header = ""
    content = ""

    with open(text_file, 'r', encoding='utf-8') as f:
        for line_num, line in enumerate(f):
            if line.startswith('#Title:'):
                header = line.replace('#Title:', '').strip()
                slide = prs.slides.add_slide(prs.slide_layouts[0])
                title = slide.shapes.title
                title.text = header
                continue
            elif line.startswith('#Slide:'):
                contents = extract_contents_from_text(f.read())
                placeholder_content = [slide for slide in contents]
                slide = prs.slides.add_slide(prs.slide_layouts[slide_count + 1])
                title = slide.shapes.title
                title.text = header
                count = 0
                for i in index_containing_placeholders[slide_count + 1]:
                    body_shape = slide.shapes.placeholders[i]
                    tf = body_shape.text_frame
                    tf.text = placeholder_content[slide_count][count]
                    count += 1
                slide_count += 1
                continue
            elif line.startswith('#Header:'):
                header = line.replace('#Header:', '').strip()
                continue

    ppt_path = f'Powerpointer-main/GeneratedPresentations/{ppt_name}.pptx'
    prs.save(ppt_path)
    return ppt_path


def create_ppt_default(text_file: str, design_number: int, ppt_name: str) -> str:
    """
    Creates a default PowerPoint presentation based on the provided text file and design number.

    Parameters:
    text_file (str): Path to the text file containing slide contents.
    design_number (int): The design number of the PowerPoint template.
    ppt_name (str): The name for the generated PowerPoint file.

    Returns:
    str: Path to the generated PowerPoint file.
    """
    prs = Presentation(f"Powerpointer-main/Designs/Design-{design_number}.pptx")
    slide_count = 0
    header = ""
    content = ""
    first_time = True

    layout_indices = [1, 7, 8]  # Valid slide layouts
    placeholder_indices = {1: 1, 7: 1, 8: 2}  # Corresponding placeholder indices

    with open(text_file, 'r', encoding='utf-8') as f:
        for line in f:
            if line.startswith('#Title:'):
                header = line.replace('#Title:', '').strip()
                slide = prs.slides.add_slide(prs.slide_layouts[0])
                title = slide.shapes.title
                title.text = header
                continue
            elif line.startswith('#Slide:'):
                slide_layout_index = random.choice(layout_indices) if not first_time else 1
                first_time = False
                slide = prs.slides.add_slide(prs.slide_layouts[slide_layout_index])
                title = slide.shapes.title
                title.text = header
                body_shape = slide.shapes.placeholders[placeholder_indices[slide_layout_index]]
                tf = body_shape.text_frame
                tf.text = content
                content = ""
                slide_count += 1
                continue
            elif line.startswith('#Header:'):
                header = line.replace('#Header:', '').strip()
                continue
            elif line.startswith('#Content:'):
                content = line.replace('#Content:', '').strip()
                next_line = f.readline().strip()
                while next_line and not next_line.startswith('#'):
                    content += '\n' + next_line
                    next_line = f.readline().strip()
                continue

    # Add the last slide if there is content remaining
    if content:
        slide_layout_index = random.choice(layout_indices) if not first_time else 1
        slide = prs.slides.add_slide(prs.slide_layouts[slide_layout_index])
        title = slide.shapes.title
        title.text = header
        body_shape = slide.shapes.placeholders[placeholder_indices[slide_layout_index]]
        tf = body_shape.text_frame
        tf.text = content

    ppt_path = f'Powerpointer-main/GeneratedPresentations/{ppt_name}.pptx'
    prs.save(ppt_path)
    return ppt_path
